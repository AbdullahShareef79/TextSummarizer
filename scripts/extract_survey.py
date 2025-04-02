import os
from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extracts text from a PowerPoint file."""
    prs = Presentation(file_path)
    text_content = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    
    return "\n".join(text_content)

def process_survey_files(input_folder, output_folder):
    """Processes all PowerPoint survey files in a folder and saves extracted text."""
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    extracted_data = {}
    for file in os.listdir(input_folder):
        if file.endswith(".pptx"):
            file_path = os.path.join(input_folder, file)
            print(f"Processing: {file}")
            text = extract_text_from_pptx(file_path)
            extracted_data[file] = text
            
            # Save extracted text to a file
            output_file = os.path.join(output_folder, f"{file}.txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)
    
    return extracted_data

if __name__ == "__main__":
    input_folder = "data/surveys"
    output_folder = "output/survey_text"
    survey_texts = process_survey_files(input_folder, output_folder)
    
    # Print extracted text summaries
    for file, text in survey_texts.items():
        print(f"\nExtracted text from {file}:\n{text[:500]}...\n")
