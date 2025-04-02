import os
from pptx import Presentation
from pptx.util import Inches

def create_summary_slide(prs, title, content):
    """Creates a summary slide with a title and content."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    content_shape.text = content

def generate_summary_pptx(summaries, output_file):
    """Generates a PowerPoint file summarizing extracted insights."""
    prs = Presentation()
    
    for category, summary in summaries.items():
        create_summary_slide(prs, category, summary)
    
    prs.save(output_file)
    print(f"Summary slides saved to {output_file}")

if __name__ == "__main__":
    summaries = {
        "Task Manager Insights": "Key findings from the survey and interviews regarding the Task Manager...",
        "Password Manager Insights": "Key findings from the survey and interviews regarding the Password Manager..."
    }
    
    output_pptx = "output/summary_presentation.pptx"
    generate_summary_pptx(summaries, output_pptx)
